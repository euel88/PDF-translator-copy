"""
PowerPoint 문서 핸들러 (.pptx, .ppt)
python-pptx 라이브러리를 사용하여 PowerPoint 문서 번역

개선된 기능:
- Run 단위 스타일 완벽 보존
- 대용량 문서 청크 처리
- 동적 배치 크기 조정
- 발표자 노트 번역
- 차트 제목/레이블 번역
- SmartArt 텍스트 번역
"""
from typing import Optional, List, Tuple, Any, Dict
from pathlib import Path
from dataclasses import dataclass

from pdf2zh.handlers.base import BaseDocumentHandler, DocumentResult


@dataclass
class RunStyle:
    """Run 스타일 정보"""
    font_name: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    font_color_rgb: Optional[Tuple[int, int, int]] = None


@dataclass
class TextLocation:
    """텍스트 위치 정보"""
    loc_type: str  # 'shape', 'table', 'note', 'chart', 'smartart'
    identifiers: tuple
    text: str
    style: Optional[RunStyle] = None


class PowerPointHandler(BaseDocumentHandler):
    """PowerPoint 문서 핸들러 - 대용량 처리 개선"""

    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']

    # 대용량 처리 설정
    BATCH_SIZE_SMALL = 40   # 짧은 텍스트
    BATCH_SIZE_MEDIUM = 25  # 중간 텍스트
    BATCH_SIZE_LARGE = 12   # 긴 텍스트

    def convert(
        self,
        input_path: str,
        output_path: Optional[str] = None
    ) -> DocumentResult:
        """
        PowerPoint 문서 번역

        Args:
            input_path: 입력 파일 경로
            output_path: 출력 파일 경로

        Returns:
            DocumentResult: 변환 결과
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            return DocumentResult(
                success=False,
                error="python-pptx 라이브러리가 설치되지 않았습니다. 'pip install python-pptx'를 실행하세요."
            )

        if output_path is None:
            output_path = self.generate_output_path(input_path)

        try:
            self.log(f"PowerPoint 문서 열기: {input_path}")
            prs = Presentation(input_path)

            # 모든 텍스트 수집 (도형, 표, 노트, 차트 등)
            text_locations = self._collect_all_texts(prs)
            total_count = len(text_locations)

            if total_count == 0:
                prs.save(output_path)
                return DocumentResult(
                    output_path=output_path,
                    success=True,
                    translated_count=0,
                    total_count=0
                )

            # 텍스트 추출
            texts_to_translate = [loc.text for loc in text_locations]

            # 동적 배치 크기로 번역
            self.log("번역 중...")
            translated_texts = self._translate_with_dynamic_batch(texts_to_translate)

            # 번역 결과 적용
            translated_count = 0
            for loc, translated in zip(text_locations, translated_texts):
                try:
                    self._apply_translation_new(prs, loc, translated)
                    translated_count += 1
                except Exception as e:
                    self.log(f"경고: 텍스트 적용 실패 ({loc.loc_type}) - {e}")

            # 저장
            self.log(f"저장 중: {output_path}")
            prs.save(output_path)

            return DocumentResult(
                output_path=output_path,
                success=True,
                translated_count=translated_count,
                total_count=total_count
            )

        except Exception as e:
            import traceback
            return DocumentResult(
                success=False,
                error=f"{str(e)}\n{traceback.format_exc()}"
            )

    def _translate_with_dynamic_batch(self, texts: List[str]) -> List[str]:
        """동적 배치 크기로 번역"""
        total = len(texts)

        # 텍스트 길이에 따라 그룹화
        short_texts = []  # < 100자
        medium_texts = []  # 100-500자
        long_texts = []  # > 500자

        for i, text in enumerate(texts):
            text_len = len(text)
            if text_len < 100:
                short_texts.append((i, text))
            elif text_len < 500:
                medium_texts.append((i, text))
            else:
                long_texts.append((i, text))

        # 결과 저장용 딕셔너리
        result_map = {}

        # 짧은 텍스트 배치 처리
        if short_texts:
            self.log(f"짧은 텍스트 번역: {len(short_texts)}개")
            self._batch_translate(short_texts, result_map, self.BATCH_SIZE_SMALL)

        # 중간 텍스트 배치 처리
        if medium_texts:
            self.log(f"중간 텍스트 번역: {len(medium_texts)}개")
            self._batch_translate(medium_texts, result_map, self.BATCH_SIZE_MEDIUM)

        # 긴 텍스트 배치 처리
        if long_texts:
            self.log(f"긴 텍스트 번역: {len(long_texts)}개")
            self._batch_translate(long_texts, result_map, self.BATCH_SIZE_LARGE)

        # 원래 순서대로 결과 정렬
        translations = []
        for i in range(total):
            translations.append(result_map.get(i, texts[i]))

        return translations

    def _batch_translate(
        self,
        indexed_texts: List[Tuple[int, str]],
        result_map: Dict[int, str],
        batch_size: int
    ):
        """배치 번역 수행"""
        for i in range(0, len(indexed_texts), batch_size):
            batch = indexed_texts[i:i + batch_size]
            texts = [t[1] for t in batch]
            indices = [t[0] for t in batch]

            try:
                translated = self.translate_texts(texts)
                for idx, trans in zip(indices, translated):
                    result_map[idx] = trans
            except Exception as e:
                self.log(f"배치 번역 오류, 개별 번역 시도: {e}")
                # 개별 번역 폴백
                for idx, text in batch:
                    try:
                        result_map[idx] = self.translate_texts([text])[0]
                    except Exception:
                        result_map[idx] = text  # 실패 시 원본 유지

            self.log(f"진행: {min(i + batch_size, len(indexed_texts))}/{len(indexed_texts)}")

    def _collect_shape_texts(
        self,
        shape,
        slide_idx: int,
        shape_idx,
        texts: List[str],
        locations: List[Tuple],
        styles: List[RunStyle]
    ):
        """도형에서 텍스트 수집 (스타일 정보 포함)"""
        # 텍스트 프레임이 있는 도형
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text and run.text.strip():
                        texts.append(run.text)
                        locations.append((
                            'shape', slide_idx, shape_idx,
                            para_idx, run_idx
                        ))
                        styles.append(self._extract_run_style(run))

        # 표
        if hasattr(shape, 'table'):
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if hasattr(cell, 'text_frame'):
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(paragraph.runs):
                                if run.text and run.text.strip():
                                    texts.append(run.text)
                                    locations.append((
                                        'table', slide_idx, shape_idx,
                                        row_idx, col_idx, para_idx, run_idx
                                    ))
                                    styles.append(self._extract_run_style(run))

        # 그룹 도형 (재귀)
        if hasattr(shape, 'shapes'):
            for sub_idx, sub_shape in enumerate(shape.shapes):
                self._collect_shape_texts(
                    sub_shape, slide_idx, f"{shape_idx}_{sub_idx}",
                    texts, locations, styles
                )

    def _extract_run_style(self, run) -> RunStyle:
        """Run에서 스타일 정보 추출"""
        font_color_rgb = None
        try:
            if run.font.color and run.font.color.rgb:
                rgb = run.font.color.rgb
                font_color_rgb = (rgb[0], rgb[1], rgb[2])
        except Exception:
            pass

        font_size = None
        try:
            if run.font.size:
                font_size = run.font.size.pt
        except Exception:
            pass

        return RunStyle(
            font_name=run.font.name,
            font_size=font_size,
            bold=run.font.bold,
            italic=run.font.italic,
            underline=run.font.underline,
            font_color_rgb=font_color_rgb
        )

    def _apply_translation(self, prs, location: Tuple, translated: str, style: RunStyle):
        """번역 결과 적용 (스타일 보존)"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        loc_type = location[0]

        if loc_type == 'shape':
            _, slide_idx, shape_idx, para_idx, run_idx = location
            slide = prs.slides[slide_idx]
            shape = self._get_shape_by_index(slide, shape_idx)
            if shape and hasattr(shape, 'text_frame'):
                run = shape.text_frame.paragraphs[para_idx].runs[run_idx]
                run.text = translated
                self._apply_style_to_run(run, style)

        elif loc_type == 'table':
            _, slide_idx, shape_idx, row_idx, col_idx, para_idx, run_idx = location
            slide = prs.slides[slide_idx]
            shape = self._get_shape_by_index(slide, shape_idx)
            if shape and hasattr(shape, 'table'):
                cell = shape.table.rows[row_idx].cells[col_idx]
                run = cell.text_frame.paragraphs[para_idx].runs[run_idx]
                run.text = translated
                self._apply_style_to_run(run, style)

    def _apply_style_to_run(self, run, style: RunStyle):
        """Run에 스타일 적용"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        try:
            if style.font_name:
                run.font.name = style.font_name
            if style.font_size:
                run.font.size = Pt(style.font_size)
            if style.bold is not None:
                run.font.bold = style.bold
            if style.italic is not None:
                run.font.italic = style.italic
            if style.underline is not None:
                run.font.underline = style.underline
            if style.font_color_rgb:
                run.font.color.rgb = RGBColor(*style.font_color_rgb)
        except Exception:
            pass  # 스타일 적용 실패해도 계속 진행

    def _get_shape_by_index(self, slide, shape_idx):
        """인덱스로 도형 가져오기 (중첩 지원)"""
        if isinstance(shape_idx, int):
            return slide.shapes[shape_idx]

        # 중첩된 경우 (예: "0_1")
        if isinstance(shape_idx, str) and '_' in shape_idx:
            indices = [int(i) for i in shape_idx.split('_')]
            shape = slide.shapes[indices[0]]
            for idx in indices[1:]:
                if hasattr(shape, 'shapes'):
                    shape = shape.shapes[idx]
                else:
                    return None
            return shape

        return slide.shapes[int(shape_idx)]

    def _collect_all_texts(self, prs) -> List[TextLocation]:
        """모든 텍스트 수집 (도형, 표, 노트, 차트, SmartArt 등)"""
        locations = []
        stats = {'shape': 0, 'table': 0, 'note': 0, 'chart': 0, 'smartart': 0}

        for slide_idx, slide in enumerate(prs.slides):
            self.log(f"슬라이드 {slide_idx + 1} 분석 중...")

            # 1. 일반 도형 및 표 텍스트 수집
            for shape_idx, shape in enumerate(slide.shapes):
                self._collect_shape_texts_new(
                    shape, slide_idx, shape_idx, locations, stats
                )

            # 2. 발표자 노트 수집
            try:
                if slide.has_notes_slide and slide.notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text and notes_frame.text.strip():
                        for para_idx, para in enumerate(notes_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                if run.text and run.text.strip():
                                    locations.append(TextLocation(
                                        loc_type='note',
                                        identifiers=(slide_idx, para_idx, run_idx),
                                        text=run.text,
                                        style=self._extract_run_style(run)
                                    ))
                                    stats['note'] += 1
            except Exception:
                pass

            # 3. 차트 텍스트 수집
            try:
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_chart:
                        chart = shape.chart
                        # 차트 제목
                        if chart.has_title and chart.chart_title:
                            title_text = self._get_chart_title_text(chart.chart_title)
                            if title_text:
                                locations.append(TextLocation(
                                    loc_type='chart_title',
                                    identifiers=(slide_idx, shape_idx, 'title'),
                                    text=title_text
                                ))
                                stats['chart'] += 1

                        # 카테고리 축 제목
                        try:
                            if hasattr(chart, 'category_axis') and chart.category_axis:
                                if hasattr(chart.category_axis, 'axis_title') and chart.category_axis.axis_title:
                                    axis_text = self._get_axis_title_text(chart.category_axis.axis_title)
                                    if axis_text:
                                        locations.append(TextLocation(
                                            loc_type='chart_axis',
                                            identifiers=(slide_idx, shape_idx, 'category_axis'),
                                            text=axis_text
                                        ))
                                        stats['chart'] += 1
                        except Exception:
                            pass

                        # 값 축 제목
                        try:
                            if hasattr(chart, 'value_axis') and chart.value_axis:
                                if hasattr(chart.value_axis, 'axis_title') and chart.value_axis.axis_title:
                                    axis_text = self._get_axis_title_text(chart.value_axis.axis_title)
                                    if axis_text:
                                        locations.append(TextLocation(
                                            loc_type='chart_axis',
                                            identifiers=(slide_idx, shape_idx, 'value_axis'),
                                            text=axis_text
                                        ))
                                        stats['chart'] += 1
                        except Exception:
                            pass

                        # 시리즈 이름 (범례)
                        try:
                            for ser_idx, series in enumerate(chart.series):
                                if hasattr(series, 'name') and series.name:
                                    if isinstance(series.name, str) and series.name.strip():
                                        locations.append(TextLocation(
                                            loc_type='chart_series',
                                            identifiers=(slide_idx, shape_idx, 'series', ser_idx),
                                            text=series.name
                                        ))
                                        stats['chart'] += 1
                        except Exception:
                            pass
            except Exception:
                pass

            # 4. SmartArt 텍스트 수집 (XML 기반)
            try:
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 15:  # MSO_SHAPE_TYPE.SMART_ART = 15
                        self._collect_smartart_text(
                            shape, slide_idx, shape_idx, locations, stats
                        )
            except Exception:
                pass

        # 통계 로깅
        total = sum(stats.values())
        self.log(f"텍스트 수집 완료 (총 {total}개): 도형 {stats['shape']}, 표 {stats['table']}, "
                 f"노트 {stats['note']}, 차트 {stats['chart']}, SmartArt {stats['smartart']}")

        return locations

    def _collect_shape_texts_new(
        self,
        shape,
        slide_idx: int,
        shape_idx,
        locations: List[TextLocation],
        stats: dict
    ):
        """도형에서 텍스트 수집 (새 방식)"""
        # 텍스트 프레임이 있는 도형
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text and run.text.strip():
                        locations.append(TextLocation(
                            loc_type='shape',
                            identifiers=(slide_idx, shape_idx, para_idx, run_idx),
                            text=run.text,
                            style=self._extract_run_style(run)
                        ))
                        stats['shape'] += 1

        # 표
        if hasattr(shape, 'table'):
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if hasattr(cell, 'text_frame'):
                        for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(paragraph.runs):
                                if run.text and run.text.strip():
                                    locations.append(TextLocation(
                                        loc_type='table',
                                        identifiers=(slide_idx, shape_idx, row_idx, col_idx, para_idx, run_idx),
                                        text=run.text,
                                        style=self._extract_run_style(run)
                                    ))
                                    stats['table'] += 1

        # 그룹 도형 (재귀)
        if hasattr(shape, 'shapes'):
            for sub_idx, sub_shape in enumerate(shape.shapes):
                self._collect_shape_texts_new(
                    sub_shape, slide_idx, f"{shape_idx}_{sub_idx}",
                    locations, stats
                )

    def _collect_smartart_text(
        self,
        shape,
        slide_idx: int,
        shape_idx: int,
        locations: List[TextLocation],
        stats: dict
    ):
        """SmartArt에서 텍스트 수집"""
        try:
            # python-pptx에서 SmartArt는 직접 접근이 제한됨
            # XML 파싱으로 텍스트 추출
            if hasattr(shape, '_element'):
                from lxml import etree
                ns = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'
                }

                # 모든 텍스트 요소 찾기
                for txt_idx, t_elem in enumerate(shape._element.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t')):
                    if t_elem.text and t_elem.text.strip():
                        locations.append(TextLocation(
                            loc_type='smartart',
                            identifiers=(slide_idx, shape_idx, txt_idx),
                            text=t_elem.text
                        ))
                        stats['smartart'] += 1
        except Exception:
            pass

    def _get_chart_title_text(self, chart_title) -> Optional[str]:
        """차트 제목에서 텍스트 추출"""
        try:
            if hasattr(chart_title, 'text_frame') and chart_title.text_frame:
                return chart_title.text_frame.text if chart_title.text_frame.text else None
            return str(chart_title) if chart_title else None
        except Exception:
            return None

    def _get_axis_title_text(self, axis_title) -> Optional[str]:
        """축 제목에서 텍스트 추출"""
        try:
            if hasattr(axis_title, 'text_frame') and axis_title.text_frame:
                return axis_title.text_frame.text if axis_title.text_frame.text else None
            return str(axis_title) if axis_title else None
        except Exception:
            return None

    def _apply_translation_new(self, prs, loc: TextLocation, translated: str):
        """번역 결과 적용 (새 방식)"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        if loc.loc_type == 'shape':
            slide_idx, shape_idx, para_idx, run_idx = loc.identifiers
            slide = prs.slides[slide_idx]
            shape = self._get_shape_by_index(slide, shape_idx)
            if shape and hasattr(shape, 'text_frame'):
                run = shape.text_frame.paragraphs[para_idx].runs[run_idx]
                run.text = translated
                if loc.style:
                    self._apply_style_to_run(run, loc.style)

        elif loc.loc_type == 'table':
            slide_idx, shape_idx, row_idx, col_idx, para_idx, run_idx = loc.identifiers
            slide = prs.slides[slide_idx]
            shape = self._get_shape_by_index(slide, shape_idx)
            if shape and hasattr(shape, 'table'):
                cell = shape.table.rows[row_idx].cells[col_idx]
                run = cell.text_frame.paragraphs[para_idx].runs[run_idx]
                run.text = translated
                if loc.style:
                    self._apply_style_to_run(run, loc.style)

        elif loc.loc_type == 'note':
            slide_idx, para_idx, run_idx = loc.identifiers
            slide = prs.slides[slide_idx]
            if slide.has_notes_slide and slide.notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    run = notes_frame.paragraphs[para_idx].runs[run_idx]
                    run.text = translated
                    if loc.style:
                        self._apply_style_to_run(run, loc.style)

        elif loc.loc_type == 'chart_title':
            slide_idx, shape_idx, _ = loc.identifiers
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if shape.has_chart and shape.chart.has_title:
                try:
                    shape.chart.chart_title.text_frame.paragraphs[0].runs[0].text = translated
                except Exception:
                    try:
                        # 대안: 전체 텍스트 설정
                        shape.chart.chart_title.text_frame.text = translated
                    except Exception:
                        pass

        elif loc.loc_type == 'chart_axis':
            slide_idx, shape_idx, axis_type = loc.identifiers
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if shape.has_chart:
                try:
                    axis = getattr(shape.chart, axis_type, None)
                    if axis and hasattr(axis, 'axis_title') and axis.axis_title:
                        if hasattr(axis.axis_title, 'text_frame'):
                            axis.axis_title.text_frame.text = translated
                except Exception:
                    pass

        elif loc.loc_type == 'chart_series':
            slide_idx, shape_idx, _, ser_idx = loc.identifiers
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            if shape.has_chart:
                try:
                    series = list(shape.chart.series)[ser_idx]
                    # 시리즈 이름은 읽기 전용일 수 있음
                    if hasattr(series, 'name'):
                        series.name = translated
                except Exception:
                    pass

        elif loc.loc_type == 'smartart':
            slide_idx, shape_idx, txt_idx = loc.identifiers
            slide = prs.slides[slide_idx]
            shape = slide.shapes[shape_idx]
            try:
                if hasattr(shape, '_element'):
                    t_elems = list(shape._element.iter(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}t'
                    ))
                    if txt_idx < len(t_elems):
                        t_elems[txt_idx].text = translated
            except Exception:
                pass
